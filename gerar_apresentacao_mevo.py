"""
Apresentacao executiva — Case Saude+ | Paleta MEVO real
1 slide = 1 mensagem | max 3 bullets | fundo claro rosa/bege
"""
import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__)) + os.sep

# ── Paleta MEVO real (extraída das imagens de marca) ────────────────────────
ROXO       = RGBColor(0x38, 0x12, 0x67)   # #381267 — roxo Mevo (cor do logo)
ROXO_MED   = RGBColor(0x5B, 0x2D, 0x8E)   # #5B2D8E — roxo médio (secundário)
ESCURO     = RGBColor(0x1A, 0x20, 0x2C)   # #1A202C — dark navy (texto)
CINZA      = RGBColor(0x71, 0x80, 0x96)   # #718096 — neutral text
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
ALERTA     = RGBColor(0xE5, 0x3E, 0x3E)   # #E53E3E — red

# ── Cores claras do site Mevo (rosa marca + bege) ───────────────────────────
ROSA_CLARO = RGBColor(0xFF, 0xE4, 0xEB)   # #FFE4EB — pink brand background
ROSA_MED   = RGBColor(0xFB, 0xBC, 0xCE)   # #FBBCCE — pink-200 (seções)
ROSA_ACENTO= RGBColor(0xF6, 0x87, 0xB3)   # #F687B3 — pink-300 (detalhe/acento)
BEGE       = RGBColor(0xFF, 0xFA, 0xF0)   # #FFFAF0 — bege quente (warmth)
CINZA_CLARO= RGBColor(0xF7, 0xFA, 0xFC)   # #F7FAFC — gray-50

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=14, bold=False, color=ESCURO,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    r.font.italic    = italic
    r.font.name      = 'Calibri'
    return tb

def img(slide, fname, x, y, w, h=None):
    path = os.path.join(BASE, fname)
    if os.path.exists(path):
        if h: slide.shapes.add_picture(path, x, y, w, h)
        else: slide.shapes.add_picture(path, x, y, w)

def logo(slide, variante='real', x=None, y=None, w=Inches(1.6)):
    """Adiciona logo Mevo real. variante: 'real' (roxo), 'branca_real', 'rosa_real'"""
    fname = f'logo_mevo_{variante}.png'
    _x = x if x is not None else W - Inches(2.0)
    _y = y if y is not None else Inches(0.12)
    img(slide, fname, _x, _y, w)

def kpi_light(slide, x, y, w, h, valor, label, accent=ROXO):
    """Card KPI estilo claro: fundo branco, borda superior colorida."""
    rect(slide, x, y, w, Inches(0.06), accent)          # borda topo
    rect(slide, x, y + Inches(0.06), w, h - Inches(0.06), BRANCO)
    txt(slide, valor,
        x + Inches(0.15), y + Inches(0.12), w - Inches(0.2), Inches(0.65),
        size=28, bold=True, color=accent, align=PP_ALIGN.LEFT)
    txt(slide, label,
        x + Inches(0.15), y + Inches(0.75), w - Inches(0.2), Inches(0.35),
        size=10, color=CINZA, align=PP_ALIGN.LEFT)

def header_light(slide, section, mensagem, cor_strip=ROXO):
    """Header claro: fundo branco, strip de cor à esq., texto escuro."""
    rect(slide, 0, 0, W, Inches(1.5), BRANCO)
    rect(slide, 0, 0, Inches(0.18), Inches(1.5), cor_strip)
    rect(slide, 0, Inches(1.5), W, Inches(0.04), RGBColor(0xE2, 0xE8, 0xF0))
    txt(slide, section,  Inches(0.38), Inches(0.1),  Inches(10), Inches(0.38),
        size=10, color=cor_strip, bold=True)
    txt(slide, mensagem, Inches(0.38), Inches(0.46), Inches(12.4), Inches(0.95),
        size=24, bold=True, color=ESCURO, wrap=True)

def bullets_light(slide, items, x, y, w, h, accent=ROXO):
    for i, item in enumerate(items[:3]):
        yi = y + i * Inches(0.75)
        rect(slide, x, yi + Inches(0.22), Inches(0.08), Inches(0.08), accent)
        txt(slide, item, x + Inches(0.2), yi, w - Inches(0.22), Inches(0.68),
            size=12.5, color=ESCURO, wrap=True)

def rosa_bg(slide):
    """Fundo rosa claro para o slide."""
    rect(slide, 0, 0, W, H, ROSA_CLARO)

def detalhe_canto(slide, cor=ROSA_MED):
    """Quadrado decorativo no canto inferior direito."""
    rect(slide, W - Inches(1.5), H - Inches(1.5), Inches(1.5), Inches(1.5), cor)
    rect(slide, W - Inches(0.8), H - Inches(0.8), Inches(0.8), Inches(0.8), BRANCO)


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, BRANCO)

# Bloco rosa no lado direito
rect(sl, Inches(8.5), 0, Inches(4.83), H, ROSA_CLARO)
rect(sl, Inches(8.5), 0, Inches(0.08), H, ROSA_MED)

# Detalhe de cor no canto sup esq
rect(sl, 0, 0, Inches(0.18), H, ROXO)

# Elementos decorativos
rect(sl, Inches(9.2), Inches(1.0), Inches(2.5), Inches(2.5), ROSA_MED)
rect(sl, Inches(10.0), Inches(1.8), Inches(2.5), Inches(2.5), BEGE)
rect(sl, Inches(9.6), Inches(1.4), Inches(2.5), Inches(2.5), BRANCO)

# Logo grande no bloco rosa
logo(sl, 'real', x=Inches(9.1), y=Inches(0.4), w=Inches(3.8))

# Tagline no bloco rosa
txt(sl, 'Mais saúde, menos complicação.',
    Inches(8.8), Inches(1.5), Inches(4.3), Inches(0.5),
    size=13, italic=True, color=CINZA, align=PP_ALIGN.CENTER)

# Texto principal (lado esq)
txt(sl, 'Analytics de Prescrições',
    Inches(0.45), Inches(1.5), Inches(7.8), Inches(0.65),
    size=36, bold=True, color=ESCURO)
txt(sl, 'Digitais',
    Inches(0.45), Inches(2.1), Inches(7.8), Inches(0.65),
    size=36, bold=True, color=ROXO)
txt(sl, 'Segmentação de Usuários & Planos de Ação',
    Inches(0.45), Inches(2.9), Inches(7.8), Inches(0.5),
    size=18, color=CINZA)

# Divisor
rect(sl, Inches(0.45), Inches(3.6), Inches(3.5), Inches(0.04), ROXO)

txt(sl, 'Janeiro – Abril 2025',
    Inches(0.45), Inches(3.8), Inches(7.8), Inches(0.4),
    size=13, color=CINZA, bold=True)
txt(sl, '70.907 prescrições  ·  68.767 pacientes  ·  15.831 médicos',
    Inches(0.45), Inches(4.2), Inches(7.8), Inches(0.4),
    size=12, color=CINZA)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — RESUMO EXECUTIVO
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rosa_bg(sl)
detalhe_canto(sl)
header_light(sl, '01  RESUMO EXECUTIVO',
             '1 em cada 2 receitas nunca chega ao paciente — e só 5% convertem digitalmente')
logo(sl, 'real', x=W - Inches(2.1), y=Inches(0.12), w=Inches(1.85))

kpis = [
    ('70.907', 'Prescrições emitidas',    ROXO),
    ('68.767', 'Pacientes únicos',        ROXO_MED),
    ('15.831', 'Médicos ativos',          ROXO_MED),
    ('50,4%',  'Open Rate geral',         RGBColor(0xD6,0x93,0x00)),
    ('10,2%',  'Conv. s/ visualizadas',   ALERTA),
    ('4.601',  'Vendas realizadas',       ROXO),
]
kw = Inches(2.1); kh = Inches(1.22); gap = Inches(0.08)
for i, (val, lbl, cor) in enumerate(kpis):
    col = i % 3; row = i // 3
    kpi_light(sl,
              Inches(0.25) + col*(kw+gap),
              Inches(1.65) + row*(kh+gap),
              kw, kh, val, lbl, accent=cor)

bullets_light(sl, [
    'Plataforma em crescimento: +22% de Jan a Mar, tendência acelerada',
    '35.135 receitas não abertas = conversão possível sem novos médicos',
    '970 compras sem rastreio digital — fechar esse gap agrega 21% ao funil',
], Inches(0.25), Inches(4.42), Inches(12.8), Inches(2.2))


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — VOLUME & SAZONALIDADE
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, BRANCO)
rect(sl, Inches(9.2), 0, Inches(4.13), H, ROSA_CLARO)
header_light(sl, '02  VOLUME & SAZONALIDADE',
             'Quinta-feira às 13h: a janela de ouro para notificar o paciente')
logo(sl, 'real', x=W - Inches(2.1), y=Inches(0.12), w=Inches(1.85))

img(sl, 'fig_mevo_volume.png', Inches(0.2), Inches(1.62), Inches(8.9), Inches(5.55))

txt(sl, 'Takeaway', Inches(9.4), Inches(1.7), Inches(3.7), Inches(0.35),
    size=9, bold=True, color=ROXO)
rect(sl, Inches(9.4), Inches(2.0), Inches(3.7), Inches(0.04), ROSA_MED)
bullets_light(sl, [
    'Pico semanal: quinta-feira (13.071 presc.) — quarta e terça também fortes',
    'Pico diário às 13h — médico emite logo após consulta matinal. SLA de notificação < 5 min',
    'Fins de semana -57% — janela ideal para re-engajar prescrições não convertidas',
], Inches(9.4), Inches(2.15), Inches(3.72), Inches(5.0))


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — PERFIL POR GERAÇÃO
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, BRANCO)
rect(sl, Inches(9.2), 0, Inches(4.13), H, BEGE)
header_light(sl, '03  PERFIL DOS PACIENTES',
             'Millennials e Gen X: 57% do volume e maior conversão — Gen Z abre mas não compra')
logo(sl, 'real', x=W - Inches(2.1), y=Inches(0.12), w=Inches(1.85))

img(sl, 'fig_mevo_geracoes.png', Inches(0.2), Inches(1.62), Inches(8.9), Inches(5.55))

txt(sl, 'Takeaway', Inches(9.4), Inches(1.7), Inches(3.7), Inches(0.35),
    size=9, bold=True, color=ROXO)
rect(sl, Inches(9.4), Inches(2.0), Inches(3.7), Inches(0.04), RGBColor(0xFE,0xD7,0xE2))
bullets_light(sl, [
    'Millennials (29-44) + Gen X (45-60) = 57% das presc. e conversão 10-13%',
    'Gen Z (13-28): OR 57,5% — o que mais abre. Mas conversão 7,5% — não compra online',
    'Boomers 70+: OR 37% — comunicar ao cuidador, não ao paciente',
], Inches(9.4), Inches(2.15), Inches(3.72), Inches(5.0))


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — OPEN RATE
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, BRANCO)
rect(sl, Inches(9.2), 0, Inches(4.13), H, ROSA_CLARO)
header_light(sl, '04  OPEN RATE',
             '35.135 receitas nunca acessadas — cada 1% de melhoria vale +709 pacientes alcançados',
             cor_strip=ALERTA)
logo(sl, 'real', x=W - Inches(2.1), y=Inches(0.12), w=Inches(1.85))

img(sl, 'fig_mevo_openrate.png', Inches(0.2), Inches(1.62), Inches(8.9), Inches(5.55))

txt(sl, 'Takeaway', Inches(9.4), Inches(1.7), Inches(3.7), Inches(0.35),
    size=9, bold=True, color=ALERTA)
rect(sl, Inches(9.4), Inches(2.0), Inches(3.7), Inches(0.04), ROSA_MED)
bullets_light(sl, [
    'OR estável em ~50% há 4 meses — sinal de estagnação sem intervenção',
    'Cirurgia Geral: OR 71,8% (topo) — investigar o que fazem diferente',
    'Meta 90d: OR 60% via push personalizado pós-emissão',
], Inches(9.4), Inches(2.15), Inches(3.72), Inches(5.0), accent=ALERTA)


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — CONVERSÃO
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, BRANCO)
rect(sl, Inches(9.2), 0, Inches(4.13), H, BEGE)
header_light(sl, '05  CONVERSÃO POR CANAL',
             '89,8% das visualizações não convertem — marketplace é o canal com maior potencial')
logo(sl, 'real', x=W - Inches(2.1), y=Inches(0.12), w=Inches(1.85))

img(sl, 'fig_mevo_conversao.png', Inches(0.2), Inches(1.62), Inches(8.9), Inches(5.55))

txt(sl, 'Takeaway', Inches(9.4), Inches(1.7), Inches(3.7), Inches(0.35),
    size=9, bold=True, color=ROXO)
rect(sl, Inches(9.4), Inches(2.0), Inches(3.7), Inches(0.04), RGBColor(0xFE,0xD7,0xE2))
bullets_light(sl, [
    'Farmácia física: 84% das vendas — offline domina, mas sem rastreio digital',
    'Marketplace: 722 vendas (16%) com potencial de 5x via redução de fricção',
    '970 compras sem abertura digital — QR Code fecha esse gap e melhora o funil',
], Inches(9.4), Inches(2.15), Inches(3.72), Inches(5.0))


# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — MÉDICOS: CONCENTRAÇÃO E CHURN
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, BRANCO)
rect(sl, Inches(9.2), 0, Inches(4.13), H, ROSA_CLARO)
header_light(sl, '06  BASE DE MÉDICOS',
             '57% dos médicos de Janeiro não prescreveram em Abril — churn é o maior risco',
             cor_strip=ALERTA)
logo(sl, 'real', x=W - Inches(2.1), y=Inches(0.12), w=Inches(1.85))

img(sl, 'fig_mevo_medicos.png', Inches(0.2), Inches(1.62), Inches(8.9), Inches(5.55))

txt(sl, 'Takeaway', Inches(9.4), Inches(1.7), Inches(3.7), Inches(0.35),
    size=9, bold=True, color=ALERTA)
rect(sl, Inches(9.4), Inches(2.0), Inches(3.7), Inches(0.04), ROSA_MED)
bullets_light(sl, [
    'VIP (92 médicos): OR 61,4% e 10,7% do volume — 1 VIP perdido = 83 presc./mês a menos',
    '36% dos médicos geram 80% das prescrições — concentração exige programa de retenção',
    'Low (9.388 med.): mediana 1 mês ativo — maioria são one-timers sem onboarding efetivo',
], Inches(9.4), Inches(2.15), Inches(3.72), Inches(5.0), accent=ALERTA)


# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — ESPECIALIDADES
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, BRANCO)
rect(sl, Inches(9.2), 0, Inches(4.13), H, BEGE)
header_light(sl, '07  ESPECIALIDADES EM DESTAQUE',
             'Psiquiatria converte 2,3x a média — medicamentos controlados criam canal obrigatório')
logo(sl, 'real', x=W - Inches(2.1), y=Inches(0.12), w=Inches(1.85))

img(sl, 'fig_mevo_especialidades.png', Inches(0.2), Inches(1.62), Inches(8.9), Inches(5.55))

txt(sl, 'Takeaway', Inches(9.4), Inches(1.7), Inches(3.7), Inches(0.35),
    size=9, bold=True, color=ROXO)
rect(sl, Inches(9.4), Inches(2.0), Inches(3.7), Inches(0.04), RGBColor(0xFE,0xD7,0xE2))
bullets_light(sl, [
    'Psiquiatria: 23,9% de conversão (2,3x a média) — 85% são controlados = canal obrigatório',
    'Cirurgia Geral: OR 71,8% mas conv. 8,9% — paciente abre mas compra no hospital',
    'Ortopedia: conv. 13% — segundo melhor; potencial com push por tipo de procedimento',
], Inches(9.4), Inches(2.15), Inches(3.72), Inches(5.0))


# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — SEGMENTAÇÃO COMPORTAMENTAL
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, BRANCO)
rect(sl, Inches(9.2), 0, Inches(4.13), H, ROSA_CLARO)
header_light(sl, '08  SEGMENTAÇÃO COMPORTAMENTAL',
             '93% das prescrições estão nos quadrantes de perda — 4 estratégias para reverter')
logo(sl, 'real', x=W - Inches(2.1), y=Inches(0.12), w=Inches(1.85))

img(sl, 'fig_mevo_comportamental.png', Inches(0.2), Inches(1.62), Inches(8.9), Inches(5.55))

# Cards de ação por quadrante no painel lateral
acoes_q = [
    (ALERTA,              'Inativo Digital  48%',   'Push < 5 min\npós-emissão'),
    (RGBColor(0xED,0x89,0x36), 'Engajado s/ Conv. 45%', 'Lembrete 2h\ne 24h'),
    (ROXO_MED,                'S/ Rastreio  1%',         'QR Code\nna receita'),
    (ROXO,               'Convertido Digital 5%',   'Renovação\nautomática'),
]
for i, (cor, titulo, acao) in enumerate(acoes_q):
    y = Inches(1.75) + i * Inches(1.38)
    rect(sl, Inches(9.35), y, Inches(3.75), Inches(1.25), BRANCO)
    rect(sl, Inches(9.35), y, Inches(0.14), Inches(1.25), cor)
    txt(sl, titulo, Inches(9.56), y + Inches(0.08), Inches(3.4), Inches(0.35),
        size=10, bold=True, color=cor)
    txt(sl, acao, Inches(9.56), y + Inches(0.48), Inches(3.4), Inches(0.65),
        size=11, color=ESCURO, wrap=True)


# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — INTELIGÊNCIA GEOGRÁFICA
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, BRANCO)
rect(sl, Inches(9.2), 0, Inches(4.13), H, BEGE)
header_light(sl, '09  INTELIGÊNCIA GEOGRÁFICA',
             'RJ e DF: alta abertura, baixa conversão — SC e RS são os modelos a replicar')
logo(sl, 'real', x=W - Inches(2.1), y=Inches(0.12), w=Inches(1.85))

img(sl, 'fig_mevo_estados.png', Inches(0.2), Inches(1.62), Inches(8.9), Inches(5.55))

txt(sl, 'Takeaway', Inches(9.4), Inches(1.7), Inches(3.7), Inches(0.35),
    size=9, bold=True, color=ROXO)
rect(sl, Inches(9.4), Inches(2.0), Inches(3.7), Inches(0.04), RGBColor(0xFE,0xD7,0xE2))
bullets_light(sl, [
    'SC: OR 61,5% + conv 12,7% — melhor combinação do país. Replicar modelo',
    'RS: conversão 14,9% (topo nacional) — investigar o que gera mais compra',
    'RJ: OR 54,7% mas conv 3,8% — anomalia crítica. Investigar estoque e preço local',
], Inches(9.4), Inches(2.15), Inches(3.72), Inches(5.0))


# ════════════════════════════════════════════════════════════════════
# SLIDE 11 — PLANOS DE AÇÃO SEGMENTADOS
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, ROSA_CLARO)
detalhe_canto(sl, ROSA_MED)
header_light(sl, '10  PLANOS DE AÇÃO SEGMENTADOS',
             '7 iniciativas com segmento, ação e meta — priorizadas por impacto')
logo(sl, 'real', x=W - Inches(2.1), y=Inches(0.12), w=Inches(1.85))

acoes_plan = [
    (ALERTA,              'Inativo\nDigital',     'Push/SMS < 5 min\npós-emissão',      'Meta:\nOR → 60%'),
    (ROXO,               'Médicos VIP\n(92)',     'Programa dedicado\nNPS mensal',       'Meta:\nChurn → 0%'),
    (ROXO_MED,                'Market-\nplace',        'Desconto 1ª\ncompra digital',         'Meta:\n20% do mix'),
    (RGBColor(0xED,0x89,0x36), 'Engajado\ns/ Conv.', 'Lembrete\n2h e 24h',             'Meta:\nConv → 13%'),
    (ROXO_MED,           'Psiquia-\ntria',        'Parceria distrib.\ncontrolados',      'Meta:\nConv → 30%'),
    (CINZA,               'RJ —\nAnomalia',        'Diagnóstico:\nestoque + UX',          'Meta:\nConv → 8%'),
    (ROSA_ACENTO,         'Crônicos\n(2,9%)',      'Renovação\nautomática',               'Meta:\nLTV +40%'),
]
cw = Inches(1.82); ch = Inches(2.45); gap = Inches(0.055)
for i, (cor, seg, acao, meta) in enumerate(acoes_plan):
    x = Inches(0.2) + i*(cw+gap)
    # header colorido
    rect(sl, x, Inches(1.6), cw, Inches(0.44), cor)
    txt(sl, seg, x, Inches(1.6), cw, Inches(0.44),
        size=9, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    # corpo branco
    rect(sl, x, Inches(2.04), cw, ch - Inches(0.44), BRANCO)
    txt(sl, 'AÇÃO', x+Inches(0.1), Inches(2.1), cw-Inches(0.12), Inches(0.26),
        size=8, bold=True, color=cor)
    txt(sl, acao, x+Inches(0.1), Inches(2.34), cw-Inches(0.12), Inches(0.9),
        size=10.5, color=ESCURO, wrap=True)
    txt(sl, 'META', x+Inches(0.1), Inches(3.28), cw-Inches(0.12), Inches(0.26),
        size=8, bold=True, color=cor)
    txt(sl, meta, x+Inches(0.1), Inches(3.52), cw-Inches(0.12), Inches(0.72),
        size=12, bold=True, color=ESCURO, wrap=True)


# ── Gerar gráfico waterfall para o PDCA ──────────────────────────────────────
def gerar_waterfall_pdca(filepath):
    fig, ax = plt.subplots(figsize=(8.5, 5.2))
    fig.patch.set_facecolor('#FFFFFF')
    ax.set_facecolor('#FFFFFF')

    itens = [
        ('Baseline\nConversao',       10.2, 'base'),
        ('Push <5 min\npos-emissao',  +0.8, 'add'),
        ('Desconto\nMarketplace',     +1.0, 'add'),
        ('Programa\nMedicos VIP',     +0.6, 'add'),
        ('Segmentacao\nde Conteudo',  +0.4, 'add'),
        ('Meta\nConversao',           13.0, 'total'),
    ]

    COR_BASE  = '#5B2D8E'
    COR_ADD   = '#FBBCCE'
    COR_TOTAL = '#381267'
    COR_TEXTO = '#1A202C'

    running = 0.0
    bottoms, vals, cores = [], [], []
    labels = [it[0] for it in itens]

    for label, val, tipo in itens:
        if tipo == 'base':
            bottoms.append(0); vals.append(val); cores.append(COR_BASE)
            running = val
        elif tipo == 'add':
            bottoms.append(running); vals.append(val); cores.append(COR_ADD)
            running += val
        else:
            bottoms.append(0); vals.append(val); cores.append(COR_TOTAL)

    x = np.arange(len(itens))
    bars = ax.bar(x, vals, bottom=bottoms, color=cores, width=0.55,
                  edgecolor='white', linewidth=1.5, zorder=3)

    for i in range(len(itens) - 2):
        top_i = bottoms[i] + vals[i]
        ax.plot([x[i] + 0.28, x[i+1] - 0.28], [top_i, top_i],
                color='#CBD5E0', lw=1.2, ls='--', zorder=2)

    for i, (bar, (label, val, tipo)) in enumerate(zip(bars, itens)):
        top = bottoms[i] + vals[i]
        if tipo == 'add':
            ax.text(x[i], top + 0.15, f'+{val:.1f} p.p.',
                    ha='center', va='bottom', fontsize=10, fontweight='bold',
                    color='#381267')
        else:
            ax.text(x[i], top + 0.15, f'{val:.1f}%',
                    ha='center', va='bottom', fontsize=11, fontweight='bold',
                    color=COR_TEXTO)

    ax.axhline(13.0, color='#381267', lw=1.5, ls=':', alpha=0.5, zorder=1)
    ax.text(5.32, 13.18, 'Meta: 13%', fontsize=9, color='#381267', va='bottom')

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=9.5, color=COR_TEXTO)
    ax.set_ylabel('Taxa de Conversao (%)', fontsize=10, color=COR_TEXTO)
    ax.set_ylim(0, 15.5)
    ax.set_xlim(-0.5, len(itens) - 0.5)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#CBD5E0')
    ax.spines['bottom'].set_color('#CBD5E0')
    ax.tick_params(colors=COR_TEXTO)
    ax.set_title('Jornada de Conversao: Baseline para Meta', fontsize=12,
                 fontweight='bold', color=COR_TEXTO, pad=10)

    patch_add   = mpatches.Patch(color=COR_ADD,   label='Iniciativa (+p.p.)')
    patch_base  = mpatches.Patch(color=COR_BASE,  label='Baseline atual')
    patch_total = mpatches.Patch(color=COR_TOTAL, label='Meta alvo')
    ax.legend(handles=[patch_base, patch_add, patch_total],
              fontsize=9, loc='upper left', framealpha=0)

    plt.tight_layout()
    plt.savefig(filepath, dpi=160, bbox_inches='tight', facecolor='white')
    plt.close()

_wf_path = os.path.join(BASE, 'fig_pdca_waterfall.png')
gerar_waterfall_pdca(_wf_path)

# ════════════════════════════════════════════════════════════════════
# SLIDE 12 — PDCA: CICLO DE MELHORIA CONTÍNUA
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, BRANCO)

# Painel esquerdo roxo (35%)
rect(sl, 0, 0, Inches(4.5), H, ROXO)
rect(sl, 0, 0, Inches(0.18), H, ROSA_ACENTO)

logo(sl, 'branca_real', x=Inches(0.35), y=Inches(0.22), w=Inches(3.1))

txt(sl, '11  PDCA', Inches(0.35), Inches(0.22), Inches(3.8), Inches(0.32),
    size=9, bold=True, color=ROSA_ACENTO)
txt(sl, 'Ciclo de\nMelhoria\nContinua',
    Inches(0.35), Inches(1.05), Inches(3.8), Inches(1.3),
    size=24, bold=True, color=BRANCO, wrap=True)

rect(sl, Inches(0.35), Inches(2.5), Inches(3.8), Inches(0.04), ROSA_MED)

pdca_bullets = [
    ('P', 'PLAN',  'OR: 60%  |  Conv: 13%  |  Churn medicos: <30%'),
    ('D', 'DO',    'Push <5 min  |  VIP program  |  Desconto marketplace  |  QR Code'),
    ('C', 'CHECK', 'OR e conversao por especialidade  |  Churn semanal  |  SLA notif.'),
    ('A', 'ACT',   'OR <55% revisar copy  |  Conv <11% testar UX  |  Churn >35% entrevistas'),
]
for i, (letra, fase, texto) in enumerate(pdca_bullets):
    y_b = Inches(2.65) + i * Inches(1.16)
    rect(sl, Inches(0.35), y_b, Inches(0.52), Inches(0.52), ROSA_MED)
    txt(sl, letra, Inches(0.35), y_b, Inches(0.52), Inches(0.52),
        size=13, bold=True, color=ROXO, align=PP_ALIGN.CENTER)
    txt(sl, fase, Inches(1.0), y_b, Inches(3.2), Inches(0.3),
        size=9, bold=True, color=ROSA_ACENTO)
    txt(sl, texto, Inches(1.0), y_b + Inches(0.28), Inches(3.3), Inches(0.85),
        size=9, color=BRANCO, wrap=True)

# Waterfall à direita
img(sl, 'fig_pdca_waterfall.png', Inches(4.55), Inches(0.35), Inches(8.6), Inches(6.9))


# ── Salvar ────────────────────────────────────────────────────────────────────
out = os.path.join(BASE, 'case_mevo_v4.pptx')
prs.save(out)
print(f'Salvo: {out}  ({len(prs.slides)} slides)')
