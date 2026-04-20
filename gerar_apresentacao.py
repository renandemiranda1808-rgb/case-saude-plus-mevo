"""
Gerador de apresentação PowerPoint — Case Técnico Analytics Saúde+
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

BASE = os.path.join(os.path.expanduser('~'), 'Projeto', 'Case Mevo')

# ── Paleta ──────────────────────────────────────────────────────────────────
AZUL       = RGBColor(0x2E, 0x86, 0xAB)
ROXO       = RGBColor(0xA2, 0x3B, 0x72)
LARANJA    = RGBColor(0xF1, 0x8F, 0x01)
VERMELHO   = RGBColor(0xC7, 0x3E, 0x1D)
VERDE      = RGBColor(0x44, 0xBB, 0xA4)
ESCURO     = RGBColor(0x1A, 0x1A, 0x2E)
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF5, 0xF5, 0xF5)
CINZA_MED  = RGBColor(0x88, 0x88, 0x88)

W = Inches(13.33)   # Widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completamente em branco

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    return shape

def add_text(slide, text, x, y, w, h, font_size=18, bold=False, color=ESCURO,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = 'Calibri'
    return txBox

def add_image(slide, path, x, y, w, h=None):
    if os.path.exists(path):
        if h:
            slide.shapes.add_picture(path, x, y, w, h)
        else:
            slide.shapes.add_picture(path, x, y, w)

def slide_header(slide, titulo, subtitulo='', cor=AZUL):
    """Faixa superior colorida com título."""
    add_rect(slide, 0, 0, W, Inches(1.2), cor)
    add_text(slide, titulo, Inches(0.4), Inches(0.15), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=BRANCO)
    if subtitulo:
        add_text(slide, subtitulo, Inches(0.4), Inches(0.82), Inches(12), Inches(0.35),
                 font_size=14, color=RGBColor(0xDD, 0xEE, 0xFF))

def add_bullet_box(slide, items, x, y, w, h, title='', title_color=AZUL, font_size=14):
    """Caixa com bullets."""
    if title:
        add_text(slide, title, x, y, w, Inches(0.35),
                 font_size=13, bold=True, color=title_color)
        y += Inches(0.38)
        h -= Inches(0.38)
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f'  •  {item}'
        p.font.size = Pt(font_size)
        p.font.color.rgb = ESCURO
        p.font.name = 'Calibri'
        p.space_after = Pt(4)

def add_kpi_card(slide, x, y, w, h, valor, label, cor):
    """Card colorido com KPI."""
    add_rect(slide, x, y, w, h, cor)
    add_text(slide, valor, x, y + Inches(0.15), w, Inches(0.65),
             font_size=32, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + Inches(0.75), w, Inches(0.4),
             font_size=12, color=BRANCO, align=PP_ALIGN.CENTER)

def add_insight_card(slide, x, y, w, h, numero, titulo, corpo, cor):
    """Card de insight com número destacado."""
    add_rect(slide, x, y, Inches(0.5), h, cor)
    add_text(slide, numero, x + Inches(0.05), y + Inches(0.05), Inches(0.4), Inches(0.4),
             font_size=18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_rect(slide, x + Inches(0.5), y, w - Inches(0.5), h, CINZA_CLARO)
    add_text(slide, titulo, x + Inches(0.6), y + Inches(0.05), w - Inches(0.7), Inches(0.35),
             font_size=13, bold=True, color=cor)
    add_text(slide, corpo, x + Inches(0.6), y + Inches(0.4), w - Inches(0.7), h - Inches(0.45),
             font_size=11, color=ESCURO, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, W, H, AZUL)
add_rect(sl, 0, Inches(3.5), W, Inches(3.5), ESCURO)

# Listras decorativas
add_rect(sl, 0, Inches(3.3), W, Inches(0.08), LARANJA)
add_rect(sl, 0, Inches(3.42), W, Inches(0.05), VERDE)

add_text(sl, 'SAÚDE+', Inches(0.6), Inches(0.6), Inches(6), Inches(0.6),
         font_size=18, bold=True, color=LARANJA)
add_text(sl, 'Case Técnico Analytics', Inches(0.6), Inches(1.3), Inches(11), Inches(1.1),
         font_size=48, bold=True, color=BRANCO)
add_text(sl, 'Análise de Prescrições Digitais — Insights & Recomendações',
         Inches(0.6), Inches(2.4), Inches(10), Inches(0.6),
         font_size=20, color=RGBColor(0xCC, 0xDD, 0xEE))
add_text(sl, 'Período: Janeiro – Abril 2025', Inches(0.6), Inches(3.85), Inches(6), Inches(0.45),
         font_size=16, color=CINZA_MED)
add_text(sl, '70.907 prescrições  |  68.767 pacientes  |  15.831 médicos',
         Inches(0.6), Inches(4.35), Inches(10), Inches(0.45),
         font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC))
add_text(sl, 'Ferramentas: Python · Pandas · Matplotlib · Seaborn',
         Inches(0.6), Inches(5.0), Inches(10), Inches(0.4),
         font_size=13, color=CINZA_MED, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, 'Agenda', cor=ESCURO)

items_esq = [
    ('01', 'Contexto & Dados',          AZUL),
    ('02', 'Resumo Executivo',          ROXO),
    ('03', 'Volume Diário & Sazonalidade', LARANJA),
    ('04', 'Perfil dos Pacientes',      AZUL),
    ('05', 'Especialidades Médicas',    VERDE),
]
items_dir = [
    ('06', 'Open Rate',                 LARANJA),
    ('07', 'Conversão por Canal',       VERMELHO),
    ('08', 'Análises Adicionais',       AZUL),
    ('09', 'Insights & Recomendações',  ROXO),
    ('10', 'Próximos Passos',           VERDE),
]

for i, (num, label, cor) in enumerate(items_esq):
    y = Inches(1.4) + i * Inches(1.0)
    add_rect(sl, Inches(0.5), y, Inches(0.55), Inches(0.6), cor)
    add_text(sl, num, Inches(0.5), y + Inches(0.1), Inches(0.55), Inches(0.45),
             font_size=18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_text(sl, label, Inches(1.15), y + Inches(0.12), Inches(5.2), Inches(0.45),
             font_size=15, color=ESCURO)

for i, (num, label, cor) in enumerate(items_dir):
    y = Inches(1.4) + i * Inches(1.0)
    add_rect(sl, Inches(6.9), y, Inches(0.55), Inches(0.6), cor)
    add_text(sl, num, Inches(6.9), y + Inches(0.1), Inches(0.55), Inches(0.45),
             font_size=18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_text(sl, label, Inches(7.55), y + Inches(0.12), Inches(5.2), Inches(0.45),
             font_size=15, color=ESCURO)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXTO & DADOS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '01  Contexto & Dados', cor=AZUL)

add_text(sl, 'Sobre a Saúde+', Inches(0.5), Inches(1.35), Inches(5.5), Inches(0.4),
         font_size=14, bold=True, color=AZUL)
add_text(sl,
    'Plataforma digital para hospitais, clínicas e consultórios médicos que integra '
    'prontuário eletrônico, emissão de prescrições com medicamentos, atestados e '
    'solicitação de exames — tudo em um só lugar.',
    Inches(0.5), Inches(1.75), Inches(5.5), Inches(1.2),
    font_size=13, color=ESCURO)

# Tabelas disponíveis
tabelas = [
    ('prescricaomedicamento.csv', '73.929 linhas', 'Prescrições com dados de paciente,\nmédico, medicamento, open rate e conversão', AZUL),
    ('medicamentos.csv',          '22 linhas',     'Detalhes dos medicamentos:\nclasse ATC, controlado, MIP, antimicrobiano', LARANJA),
    ('medicos.csv',               '45.025 linhas', 'Médicos com especialidade,\nestado e conselho profissional', ROXO),
]
for i, (nome, qtd, desc, cor) in enumerate(tabelas):
    y = Inches(1.35) + i * Inches(1.55)
    x = Inches(6.4)
    add_rect(sl, x, y, Inches(6.4), Inches(1.35), cor)
    add_text(sl, nome, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.4),
             font_size=13, bold=True, color=BRANCO)
    add_text(sl, qtd, x + Inches(0.2), y + Inches(0.48), Inches(5.5), Inches(0.3),
             font_size=20, bold=True, color=BRANCO)
    add_text(sl, desc, x + Inches(0.2), y + Inches(0.82), Inches(5.8), Inches(0.5),
             font_size=10, color=RGBColor(0xEE, 0xEE, 0xFF))

# Período
add_rect(sl, Inches(0.5), Inches(5.2), Inches(5.5), Inches(0.7), CINZA_CLARO)
add_text(sl, '📅  Período analisado: 01/01/2025  →  18/04/2025  (108 dias)',
         Inches(0.7), Inches(5.3), Inches(5.2), Inches(0.5),
         font_size=13, bold=True, color=ESCURO)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RESUMO EXECUTIVO (KPIs)
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '02  Resumo Executivo', 'Principais métricas do período Jan–Abr 2025', cor=ROXO)

kpis = [
    ('70.907',  'Prescrições Emitidas',  AZUL),
    ('68.767',  'Pacientes Únicos',      ROXO),
    ('15.831',  'Médicos Ativos',        VERDE),
    ('50,4%',   'Open Rate',             LARANJA),
    ('10,2%',   'Conv. s/ Visualizadas', VERMELHO),
    ('4.601',   'Vendas Realizadas',     RGBColor(0x3B, 0x1F, 0x2B)),
]
cols = 3
card_w = Inches(4.0)
card_h = Inches(1.45)
gap_x  = Inches(0.22)
gap_y  = Inches(0.25)
start_x = Inches(0.4)
start_y = Inches(1.45)

for i, (val, lbl, cor) in enumerate(kpis):
    col = i % cols
    row = i // cols
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_kpi_card(sl, x, y, card_w, card_h, val, lbl, cor)

add_text(sl,
    '* Conversão de 10,2% calculada sobre prescrições visualizadas (35.772). Total de vendas: 4.601 (incl. 970 via farmácia física sem visualização digital prévia).',
    Inches(0.4), Inches(6.9), Inches(12), Inches(0.35),
    font_size=10, color=CINZA_MED, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Q1: VOLUME DIÁRIO & SAZONALIDADE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '03  Volume Diário & Sazonalidade', 'Q1 — Quantas prescrições foram emitidas por dia?', cor=LARANJA)

img_path = os.path.join(BASE, 'fig_q1_volume_diario.png')
add_image(sl, img_path, Inches(0.25), Inches(1.3), Inches(8.8), Inches(5.7))

add_bullet_box(sl, [
    'Média de 657 prescrições/dia no período',
    'Pico de 1.271 prescrições em 15/04/2025',
    'Crescimento mês a mês: Jan 17k → Mar 21k (+22%)',
    'Queda expressiva nos fins de semana (Sáb/Dom): -57% vs dias úteis',
    'Pico semanal na Quinta-feira (13.071 presc.) — Qua/Ter também fortes',
    'Tendência de crescimento evidenciada pela\nmédia móvel de 7 dias',
], Inches(9.2), Inches(1.4), Inches(3.9), Inches(5.5),
   title='Destaques', title_color=LARANJA, font_size=13)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Q2: PERFIL DOS PACIENTES
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '04  Perfil dos Pacientes', 'Q2 — Quantos pacientes foram atendidos?', cor=AZUL)

img_path = os.path.join(BASE, 'fig_q2_pacientes.png')
add_image(sl, img_path, Inches(0.25), Inches(1.3), Inches(8.8), Inches(4.5))

add_kpi_card(sl, Inches(9.2), Inches(1.4), Inches(3.9), Inches(1.1),
             '68.767', 'Pacientes Únicos', AZUL)

add_bullet_box(sl, [
    '97,1% dos pacientes tiveram\napenas 1 prescrição no período',
    '2,9% retornaram (tratamento crônico)',
    'Maioria feminina (~60%)',
    'Faixa etária concentrada: 30–59 anos',
    'SP representa ~70% do volume',
    'Média de 1,03 prescrições por paciente',
], Inches(9.2), Inches(2.65), Inches(3.9), Inches(4.2),
   title='Destaques', title_color=AZUL, font_size=12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Q3: ESPECIALIDADES MÉDICAS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '05  Especialidades Médicas', 'Q3 — Volume de prescrições e produtividade por especialidade', cor=VERDE)

# Gráfico principal (volume + médicos únicos): largura total, topo
img_path = os.path.join(BASE, 'fig_q3_especialidades.png')
add_image(sl, img_path, Inches(0.2), Inches(1.3), Inches(13.0), Inches(3.3))

# Gráfico de produtividade: metade esquerda inferior
img2 = os.path.join(BASE, 'fig_q3_produtividade.png')
add_image(sl, img2, Inches(0.2), Inches(4.7), Inches(7.5), Inches(2.6))

# Insights: metade direita inferior
add_bullet_box(sl, [
    'Clínica Médica lidera: 15.293 presc. (21,6%)',
    '"Sem Especialidade": 12.959 presc. (18,3%)',
    '  96% são CRM sem especialidade preenchida —',
    '  não CROs. Risco de cadastro incompleto.',
    'Pediatria: 8.702 presc. (12,3%)',
    'Top 3 concentram 52,2% do volume total',
    'Endocrinologia: maior produtividade (10,9 presc/médico)',
    '263 prescrições sem match no cadastro (0,4%)',
], Inches(7.9), Inches(4.7), Inches(5.3), Inches(2.6),
   title='Destaques', title_color=VERDE, font_size=11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Q4: OPEN RATE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '06  Open Rate', 'Q4 — Taxa de abertura da receita digital', cor=LARANJA)

# KPI destaque
add_kpi_card(sl, Inches(0.4), Inches(1.35), Inches(2.8), Inches(1.4),
             '50,4%', 'Open Rate Geral', LARANJA)
add_kpi_card(sl, Inches(3.4), Inches(1.35), Inches(2.8), Inches(1.4),
             '35.772', 'Prescrições Visualizadas', AZUL)
add_kpi_card(sl, Inches(6.4), Inches(1.35), Inches(2.8), Inches(1.4),
             '35.135', 'Não Visualizadas', VERMELHO)

img_path = os.path.join(BASE, 'fig_q4_openrate.png')
add_image(sl, img_path, Inches(0.25), Inches(2.9), Inches(6.5), Inches(4.2))

img2 = os.path.join(BASE, 'fig_q4_openrate_esp.png')
add_image(sl, img2, Inches(6.8), Inches(2.9), Inches(6.3), Inches(4.2))

add_text(sl,
    'Definição:  Open Rate = prescrições visualizadas / total de prescrições emitidas',
    Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3),
    font_size=10, color=CINZA_MED, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Q5: CONVERSÃO POR CANAL
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '07  Conversão por Canal', 'Q5 — Taxa de conversão: prescrições vendidas / visualizadas', cor=VERMELHO)

add_kpi_card(sl, Inches(0.4),  Inches(1.35), Inches(3.0), Inches(1.3),
             '10,2%', 'Conversão s/ Visualizadas', VERMELHO)
add_kpi_card(sl, Inches(3.6),  Inches(1.35), Inches(3.0), Inches(1.3),
             '3.879', 'Vendas — Farmácia Física', AZUL)
add_kpi_card(sl, Inches(6.8),  Inches(1.35), Inches(3.0), Inches(1.3),
             '722',   'Vendas — Marketplace', VERDE)
add_kpi_card(sl, Inches(10.0), Inches(1.35), Inches(3.0), Inches(1.3),
             '6,5%',  'Conversão Geral', LARANJA)

img_path = os.path.join(BASE, 'fig_q5_conversao.png')
add_image(sl, img_path, Inches(0.25), Inches(2.8), Inches(6.5), Inches(4.0))

img2 = os.path.join(BASE, 'fig_q5_conv_esp.png')
add_image(sl, img2, Inches(6.8), Inches(2.8), Inches(6.3), Inches(4.0))

add_text(sl,
    'Farmácia física: 84% das vendas (3.879)  |  Marketplace: 16% (722)  |  Total: 4.601 vendas  |  Conv. 10,2% calculada sobre 35.772 visualizadas',
    Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3),
    font_size=10, color=CINZA_MED, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Q6: ANÁLISES ADICIONAIS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '08  Análises Adicionais', 'Q6 — Outras informações relevantes dos dados', cor=AZUL)

img1 = os.path.join(BASE, 'fig_q6_pareto.png')
add_image(sl, img1, Inches(0.25), Inches(1.35), Inches(6.2), Inches(3.0))

img2 = os.path.join(BASE, 'fig_q6_hora.png')
add_image(sl, img2, Inches(0.25), Inches(4.45), Inches(6.2), Inches(2.7))

add_bullet_box(sl, [
    'Pareto: 36,4% dos médicos geram 80% das prescrições',
    'Top 10 médicos = 2,5% do volume — oportunidade VIP',
    '5.767 médicos-chave para programa de parceiros',
    'Pico de emissão: 13h — janela crítica para notificação',
    'Manhã (9h–11h) e tarde (13h–17h) = 78% do volume diário',
    'Medicamento mais prescrito: Paracetamol (12.895 itens)',
    '2º Glifage XR (8.997) | 3º Aerolin (6.516) | 4º Avamys (6.271)',
    'Tempo mediano até conversão: 8,2h',
    '64% das compras nas primeiras 24h — janela decisiva',
], Inches(6.6), Inches(1.4), Inches(6.5), Inches(5.7),
   title='Destaques', title_color=AZUL, font_size=12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Q7: INSIGHTS 1, 2, 3
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '09  Insights & Recomendações  (1/2)', 'Q7 — Insights acionáveis para melhorar o produto', cor=ROXO)

insights_a = [
    ('1', 'Open Rate de 50% — metade das receitas não é aberta',
     'Implementar notificações push/SMS imediatamente após emissão. '
     'Testar mensagens personalizadas com nome do medicamento. '
     'Monitorar impacto por faixa etária e especialidade.',
     LARANJA),
    ('2', 'Sazonalidade semanal: quedas nos fins de semana',
     'Usar sábado/domingo para re-engajamento: enviar lembrete a pacientes que '
     'visualizaram mas não compraram. Planejar capacidade técnica e de suporte '
     'conforme picos nos dias úteis.',
     AZUL),
    ('3', 'Concentração de médicos: top 36% geram 80% das prescrições',
     'Criar programa "Médicos Parceiros" (key accounts) com suporte prioritário '
     'e treinamento. Investigar abandono dos médicos de baixo volume: '
     'dificuldade técnica, especialidade inadequada ou churn.',
     VERDE),
]

for i, (num, titulo, corpo, cor) in enumerate(insights_a):
    y = Inches(1.4) + i * Inches(1.85)
    add_insight_card(sl, Inches(0.4), y, Inches(12.5), Inches(1.65), num, titulo, corpo, cor)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Q7: INSIGHTS 4, 5, 6, 7
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '09  Insights & Recomendações  (2/2)', cor=ROXO)

insights_b = [
    ('4', 'Marketplace: canal com maior potencial de escala digital',
     'Investigar fricção no marketplace (estoque, preço, UX). Testar campanha '
     'de desconto na 1ª compra digital. Identificar especialidades com alta '
     'conversão e baixa penetração no canal.',
     VERMELHO),
    ('5', 'Pacientes crônicos: 2,9% retornaram com 2+ prescrições',
     'Criar funcionalidade de renovação automática para uso contínuo (Glifage, '
     'Lítio, Fluoxetina). Lembrete proativo de renovação. Crônicos têm alto LTV.',
     ROXO),
    ('6', 'Pico de prescrições às 13h (início da tarde)',
     'Garantir SLA de notificação < 5 min no horário de pico (13h). '
     'Segundos picos: 9h–11h (manhã) e 15h–17h (tarde). '
     'Alocar equipe de suporte alinhada com a curva de demanda.',
     LARANJA),
    ('7', 'CROs cadastrados como "Sem Especialidade" — risco de compliance',
     'Enriquecer cadastro de médicos para separar CRM e CRO. '
     'Auditar prescrições de CROs fora do escopo odontológico. '
     'Criar relatórios de prescrição por tipo de profissional.',
     VERDE),
    ('8', '970 vendas em farmácia física sem visualização digital — oportunidade de rastreabilidade',
     '21% das vendas em farmácia física ocorreram sem abertura da receita digital. '
     'Investigar se o canal offline captura prescrições por outros meios. '
     'Implementar tracking unificado para fechar o funil de conversão real.',
     AZUL),
]

for i, (num, titulo, corpo, cor) in enumerate(insights_b):
    y = Inches(1.4) + i * Inches(1.16)
    add_insight_card(sl, Inches(0.4), y, Inches(12.5), Inches(1.0), num, titulo, corpo, cor)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — PDCA: PRÓXIMOS PASSOS & MÉTRICAS DE ACOMPANHAMENTO
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '10  Ciclo PDCA — Próximos Passos & Métricas', 'Da análise à melhoria contínua', cor=VERDE)

# ── Quatro quadrantes PDCA ──────────────────────────────────────────────────
pdca = [
    (AZUL,    'P  PLAN — Planejar',
     '• Meta Open Rate: 60% (+9,6 p.p. em 90 dias)\n'
     '• Meta Conversão: 13% (+2,8 p.p. em 90 dias)\n'
     '• Meta Marketplace: 20% do mix de vendas\n'
     '• Programa Médicos Parceiros: top 5.767\n'
     '• Enriquecer cadastro: zerar "Sem Especialidade"'),
    (LARANJA, 'D  DO — Executar',
     '• Push/SMS < 5 min após emissão (13h e 9h–11h)\n'
     '• Re-engajamento fim de semana (não convertidos)\n'
     '• Incentivo marketplace: desconto 1ª compra\n'
     '• Funcionalidade renovação de receita (crônicos)\n'
     '• Rastreio offline: QR Code farmácia física'),
    (VERDE,   'C  CHECK — Verificar',
     '• Open Rate: semana a semana por especialidade\n'
     '• Conversão: funil digital vs físico por canal\n'
     '• Churn de médicos: ativos mês anterior vs atual\n'
     '• Tempo médio notificação: SLA < 5 min (% atingido)\n'
     '• Vendas marketplace: % do mix semanal'),
    (ROXO,    'A  ACT — Agir',
     '• Open Rate < 55%: revisar copy e horário do push\n'
     '• Conv. < 11%: testar desconto ou simplificar checkout\n'
     '• Churn médico > 5%: NPS + entrevistas de cancelamento\n'
     '• Marketplace < 18%: investigar estoque e UX\n'
     '• Revisão PDCA: quinzenal com time de produto'),
]

for i, (cor, titulo, corpo) in enumerate(pdca):
    col = i % 2
    row = i // 2
    x = Inches(0.3) + col * Inches(6.6)
    y = Inches(1.4) + row * Inches(2.85)
    w = Inches(6.4)
    h = Inches(2.65)
    add_rect(sl, x, y, Inches(0.18), h, cor)
    add_rect(sl, x + Inches(0.18), y, w - Inches(0.18), h, CINZA_CLARO)
    add_text(sl, titulo, x + Inches(0.28), y + Inches(0.08), w - Inches(0.35), Inches(0.38),
             font_size=13, bold=True, color=cor)
    add_text(sl, corpo, x + Inches(0.28), y + Inches(0.48), w - Inches(0.35), h - Inches(0.55),
             font_size=10.5, color=ESCURO, wrap=True)

# ── Linha de métricas de controle na base ───────────────────────────────────
add_rect(sl, 0, Inches(7.0), W, Inches(0.5), ESCURO)
metricas = [
    ('Open Rate',  '50,4%', '→ Meta 60%'),
    ('Conversão',  '10,2%', '→ Meta 13%'),
    ('Marketplace','16%',   '→ Meta 20%'),
    ('SLA Notif.', '–',     '→ Meta <5min'),
    ('Churn Med.', '–',     '→ Meta <5%/mês'),
]
for i, (label, atual, meta) in enumerate(metricas):
    x = Inches(0.5) + i * Inches(2.6)
    add_text(sl, f'{label}', x, Inches(7.03), Inches(2.5), Inches(0.2),
             font_size=8, bold=True, color=LARANJA)
    add_text(sl, f'{atual}  {meta}', x, Inches(7.22), Inches(2.5), Inches(0.22),
             font_size=9, color=BRANCO)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — SEGMENTAÇÃO COMPORTAMENTAL
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '11  Segmentação Comportamental', 'Matriz de Engajamento × Conversão — 70.907 prescrições', cor=VERMELHO)

add_image(sl, os.path.join(BASE, 'fig_seg_comportamental.png'),
          Inches(0.2), Inches(1.3), Inches(13.0), Inches(4.3))

# 4 cards de ação por segmento
acoes_seg = [
    (VERMELHO, 'Inativo Digital — 48,2%',
     'Presc. nunca abertas. Acao:\npush/SMS < 5 min pos-emissao.\nMeta: reduzir para 40%.'),
    (LARANJA,  'Engajado Nao Convertido — 45,3%',
     'Abriu mas nao comprou. Acao:\nlembrete 2h e 24h pos-abertura.\nOportunidade: +970 vendas/mes.'),
    (AZUL,     'Comprou Sem Digital — 1,4%',
     'Farmacia fisica sem rastreio.\nAcao: QR Code na receita.\nMeta: mapear canal offline.'),
    (VERDE,    'Convertido Digital — 5,1%',
     'Abriu e comprou. Acao:\nfidelizar com renovacao automatica.\nMeta: 7% em 90 dias.'),
]
card_w = Inches(3.2)
for i, (cor, titulo, corpo) in enumerate(acoes_seg):
    x = Inches(0.2) + i * Inches(3.28)
    add_rect(sl, x, Inches(5.7), card_w, Inches(1.6), cor)
    add_text(sl, titulo, x + Inches(0.12), Inches(5.75), card_w - Inches(0.15), Inches(0.38),
             font_size=10, bold=True, color=BRANCO)
    add_text(sl, corpo,  x + Inches(0.12), Inches(6.14), card_w - Inches(0.15), Inches(1.1),
             font_size=9, color=BRANCO, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — SEGMENTAÇÃO DE MÉDICOS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '12  Segmentação de Médicos', 'Volume × Engajamento × Retenção — 15.831 médicos ativos', cor=ROXO)

add_image(sl, os.path.join(BASE, 'fig_seg_medicos.png'),
          Inches(0.2), Inches(1.3), Inches(13.0), Inches(3.8))

# KPI de churn em destaque
add_rect(sl, Inches(0.2), Inches(5.2), Inches(4.0), Inches(2.1), VERMELHO)
add_text(sl, '57%', Inches(0.2), Inches(5.3), Inches(4.0), Inches(0.9),
         font_size=42, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_text(sl, 'Churn Jan→Abr\n4.105 medicos nao voltaram',
         Inches(0.2), Inches(6.1), Inches(4.0), Inches(0.7),
         font_size=10, color=BRANCO, align=PP_ALIGN.CENTER)

planos_med = [
    (ROXO,  'VIP (92 med.) — 10,7% do volume',
     'OR 61,4% — melhor segmento. Programa exclusivo:\nSuporte dedicado, NPS mensal, acesso antecipado a funcionalidades. Meta: 0% churn.'),
    (AZUL,  'High (1.596 med.) — 40% do volume',
     'OR 50%. Foco em ativacao de prescrições. Treinamento, webinar mensal e feedback de produto. Meta: migrar 200 para VIP.'),
    (LARANJA,'Mid (4.755 med.) — 32% do volume',
     'OR 47%. Medicos em crescimento. Onboarding automatizado e notificacoes de desempenho. Meta: converter 500 para High.'),
    (VERDE,  'Low (9.388 med.) — 17% do volume',
     'OR 49% mas mediana de 1 mes ativo: sao mostly one-timers. Investigar causa de abandono. Reativacao por email 30d apos inatividade.'),
]
for i, (cor, titulo, corpo) in enumerate(planos_med):
    x = Inches(4.4) + i * Inches(2.23)
    add_rect(sl, x, Inches(5.2), Inches(2.1), Inches(2.1), cor)
    add_text(sl, titulo, x + Inches(0.08), Inches(5.25), Inches(1.95), Inches(0.4),
             font_size=8, bold=True, color=BRANCO, wrap=True)
    add_text(sl, corpo, x + Inches(0.08), Inches(5.65), Inches(1.95), Inches(1.55),
             font_size=8, color=BRANCO, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — SEGMENTAÇÃO POR PERFIL
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, '13  Segmentação por Perfil', 'Faixa Etária × Estado × Especialidade — Open Rate vs Conversão', cor=AZUL)

add_image(sl, os.path.join(BASE, 'fig_seg_perfil.png'),
          Inches(0.2), Inches(1.3), Inches(13.0), Inches(3.8))

# Destaques e planos de ação segmentados
destaques = [
    (AZUL,   'Faixa 30–49 anos\nMelhor perfil de conversao',
     'OR ~55% e conv 10-13%. Prioridade para campanhas de marketplace e renovacao de receita. Maior LTV esperado.'),
    (VERDE,  'SC, RS, SP\nTop performers geograficos',
     'Conv 12-15%, acima da media. Estrategia: ampliar cobertura de medicos e farmácias parceiras nesses estados.'),
    (ROXO,   'Psiquiatria — conv 23,9%',
     'Maior conversao dentre todas as especialidades. 85% prescricoes controladas. Parceria com distribuidoras especializadas + marketplace para controlados.'),
    (LARANJA,'RJ — anomalia critica',
     'OR 54,7% mas conv 3,8% (vs media 10,2%). Investigar: estoque na farmacia fisica, concorrencia local, UX do marketplace no RJ.'),
    (VERMELHO,'0–17 e 70+',
     'OR e conv abaixo da media. Publico dependente de responsavel. Comunicacao adaptada ao cuidador, nao ao paciente.'),
    (CINZA_MED, 'Cirurgia Geral — OR 71,8%',
     'Maior OR, mas conv baixa (8,9%). Paciente abre mas nao compra. Investigar se medicamentos sao comprados no proprio hospital/clinica.'),
]
for i, (cor, titulo, corpo) in enumerate(destaques):
    col = i % 3
    row = i // 3
    x = Inches(0.2) + col * Inches(4.38)
    y = Inches(5.2) + row * Inches(1.05)
    add_rect(sl, x, y, Inches(0.12), Inches(0.95), cor)
    add_rect(sl, x + Inches(0.12), y, Inches(4.12), Inches(0.95), RGBColor(0xF5, 0xF5, 0xF5))
    add_text(sl, titulo, x + Inches(0.2), y + Inches(0.04), Inches(4.0), Inches(0.3),
             font_size=9, bold=True, color=cor)
    add_text(sl, corpo, x + Inches(0.2), y + Inches(0.34), Inches(4.0), Inches(0.58),
             font_size=8, color=ESCURO, wrap=True)


# ── Salvar ───────────────────────────────────────────────────────────────────
out = os.path.join(BASE, 'case_saude_plus_v2.pptx')
prs.save(out)
print(f'Apresentacao gerada: {out}')
print(f'   {len(prs.slides)} slides criados')
